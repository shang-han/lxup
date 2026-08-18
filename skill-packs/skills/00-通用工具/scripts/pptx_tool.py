#!/usr/bin/env python3
# -*- coding: utf-8 -*-
try:
    import sys
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")
except Exception:
    pass

import argparse, sys
try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    sys.stderr.write('[pptx_tool] 缺少依赖 python-pptx,请先执行: pip install "python-pptx>=0.6.21"\n'); sys.exit(1)

CJK_FONT = '微软雅黑'

def is_cjk(s):
    return any('一' <= ch <= '鿿' for ch in s)

def set_run(para, text, size, lvl=0):
    para.level = lvl
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    if is_cjk(text):
        run.font.name = CJK_FONT  # 中文字体自适应

def parse_outline(path):
    """解析大纲:# 封面 / ## 页标题 / - 一级要点 / '  - ' 二级要点 / > 讲解备注"""
    try:
        lines = open(path, encoding='utf-8').read().splitlines()
    except Exception as e:
        sys.stderr.write('[pptx_tool] 无法读取大纲: ' + str(e) + '\n'); sys.exit(1)
    cover_title = None
    slides = []
    cur = None
    for ln in lines:
        raw = ln.rstrip()
        if raw.startswith('# '):
            cover_title = raw[2:].strip()
        elif raw.startswith('## '):
            cur = {'title': raw[3:].strip(), 'items': [], 'notes': []}
            slides.append(cur)
        elif cur is not None:
            if raw.startswith('- '):
                cur['items'].append((0, raw[2:].strip()))
            elif raw.startswith('  - '):
                cur['items'].append((1, raw[4:].strip()))
            elif raw.startswith('> '):
                cur['notes'].append(raw[2:].strip())
    return cover_title, slides

def build_content_slide(prs, item):
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    sl.shapes.title.text = item['title']
    if is_cjk(item['title']):
        for para in sl.shapes.title.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = CJK_FONT
    tf = sl.placeholders[1].text_frame
    tf.clear()
    n = len(item['items'])
    # 要点超过 5 条自动缩字号,保证不溢出
    size0, size1 = (18, 16) if n > 5 else (24, 20)
    for i, (lvl, text) in enumerate(item['items']):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        set_run(para, text, size1 if lvl else size0, lvl)
    if item['notes']:
        sl.notes_slide.notes_text_frame.text = '\n'.join(item['notes'])
    return sl

def main():
    p = argparse.ArgumentParser()
    p.add_argument('create'); p.add_argument('-o','--output', required=True); p.add_argument('--outline', required=True)
    a = p.parse_args()
    cover_title, slides = parse_outline(a.outline)
    if not slides and not cover_title:
        sys.stderr.write('[pptx_tool] 大纲无内容(需 # 或 ## 开头)\n'); sys.exit(1)
    prs = Presentation()
    if cover_title:
        sl = prs.slides.add_slide(prs.slide_layouts[0])
        sl.shapes.title.text = cover_title
        if is_cjk(cover_title):
            for para in sl.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = CJK_FONT
    for item in slides:
        build_content_slide(prs, item)
    prs.save(a.output)
    print('已生成 ' + a.output + ',共 ' + str(len(prs.slides._sldIdLst)) + ' 页')

main()
